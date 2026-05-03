#!/usr/bin/env python3
"""
document_generator.py - JARVIS Document Generator Module
=========================================================

Generates professional documents programmatically:
  - PDFs (with Hebrew font support)
  - Word documents (.docx)
  - PowerPoint presentations (.pptx)
  - Excel spreadsheets (.xlsx)
  - HTML and Markdown

Specialized templates:
  - Legal contracts (NDA, employment, service, partnership, license)
  - Medical reports
  - Technical specifications
  - Business plans
  - Meeting minutes
  - Project proposals

Supports Hebrew and English. 100% local generation.

Dependencies (graceful degradation):
  - reportlab  → PDF generation
  - python-docx → Word documents
  - python-pptx → Presentations
  - openpyxl   → Spreadsheets
  - fpdf       → PDF fallback
"""

from __future__ import annotations

import os
import subprocess
import sys
import warnings

# ---------------------------------------------------------------------------
# Logging: the local logging.py shadows stdlib logging, so we roll our own
# lightweight logger that forwards to warnings/showwarning.
# ---------------------------------------------------------------------------

class _LightweightLogger:
    """Minimal logger compatible with logging.Logger interface."""
    _level_map = {"debug": 10, "info": 20, "warning": 30, "error": 40, "critical": 50}

    def __init__(self, name: str, level: int = 20):
        self.name = name
        self.level = level

    def isEnabledFor(self, level: int) -> bool:
        return level >= self.level

    def debug(self, msg: str, *args) -> None:
        if self.isEnabledFor(10):
            print(f"[DEBUG] {self.name}: {msg % args if args else msg}")

    def info(self, msg: str, *args) -> None:
        if self.isEnabledFor(20):
            print(f"[INFO]  {self.name}: {msg % args if args else msg}")

    def warning(self, msg: str, *args) -> None:
        if self.isEnabledFor(30):
            print(f"[WARN]  {self.name}: {msg % args if args else msg}")

    def error(self, msg: str, *args) -> None:
        if self.isEnabledFor(40):
            print(f"[ERROR] {self.name}: {msg % args if args else msg}")

    def critical(self, msg: str, *args) -> None:
        if self.isEnabledFor(50):
            print(f"[CRIT]  {self.name}: {msg % args if args else msg}")

    def exception(self, msg: str, *args) -> None:
        import traceback
        self.error(msg, *args)
        traceback.print_exc()


# Create lightweight logger at module level (avoids stdlib logging import issues)
logger = _LightweightLogger("document_generator")
import tempfile
import textwrap
import time
from datetime import datetime
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

# ---------------------------------------------------------------------------
# Optional third-party imports with graceful fallbacks
# ---------------------------------------------------------------------------
try:
    from reportlab.lib import colors as rl_colors
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY
    from reportlab.lib.pagesizes import A4, letter
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import cm, mm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import (
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
        PageBreak,
        ListFlowable,
        ListItem,
    )
    HAS_REPORTLAB = True
except Exception:
    HAS_REPORTLAB = False

try:
    from docx import Document
    from docx.shared import Inches, Pt, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.style import WD_STYLE_TYPE
    HAS_PYTHON_DOCX = True
except Exception:
    HAS_PYTHON_DOCX = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor as PptxRGBColor
    HAS_PYTHON_PPTX = True
except Exception:
    HAS_PYTHON_PPTX = False

try:
    import openpyxl
    from openpyxl.styles import (
        Font,
        Alignment,
        PatternFill,
        Border,
        Side,
        NamedStyle,
    )
    HAS_OPENPYXL = True
except Exception:
    HAS_OPENPYXL = False

try:
    from fpdf import FPDF
    HAS_FPDF = True
except Exception:
    HAS_FPDF = False

# ============================================================================
# Constants
# ============================================================================

HEBREW_FONT_FAMILIES = [
    "DejaVuSans",
    "Arial",
    "Times New Roman",
    "Courier New",
]

HEBREW_MONTHS = {
    "he": [
        "ינואר", "פברואר", "מרץ", "אפריל", "מאי", "יוני",
        "יולי", "אוגוסט", "ספטמבר", "אוקטובר", "נובמבר", "דצמבר",
    ],
    "en": [
        "January", "February", "March", "April", "May", "June",
        "July", "August", "September", "October", "November", "December",
    ],
}

LEGAL_TERMS_HE = {
    "parties": "הצדדים",
    "preamble": "הקדמה",
    "terms": "תנאים והוראות",
    "confidentiality": "סודיות",
    "termination": "סיום",
    "governing_law": "דין חל",
    "signatures": "חתימות",
    "witness": "עד",
    "agreement": "הסכם",
    "effective_date": "תאריך תחילת תוקף",
    "obligation": "התחייבות",
    "liability": "אחריות",
    "indemnification": "פיצוי",
}

LEGAL_TERMS_EN = {
    "parties": "Parties",
    "preamble": "Preamble",
    "terms": "Terms and Conditions",
    "confidentiality": "Confidentiality",
    "termination": "Termination",
    "governing_law": "Governing Law",
    "signatures": "Signatures",
    "witness": "Witness",
    "agreement": "Agreement",
    "effective_date": "Effective Date",
    "obligation": "Obligation",
    "liability": "Liability",
    "indemnification": "Indemnification",
}

MEDICAL_SECTIONS_HE = [
    ("פרטי מטופל", "patient_info"),
    ("סיבת הפניה", "referral_reason"),
    ("תלונות עיקריות", "chief_complaints"),
    ("בדיקה פיזיקלית", "physical_examination"),
    ("ממצאים", "findings"),
    ("אבחנה", "diagnosis"),
    ("המלצות", "recommendations"),
    ("מעקב", "follow_up"),
]

MEDICAL_SECTIONS_EN = [
    ("Patient Information", "patient_info"),
    ("Referral Reason", "referral_reason"),
    ("Chief Complaints", "chief_complaints"),
    ("Physical Examination", "physical_examination"),
    ("Findings", "findings"),
    ("Diagnosis", "diagnosis"),
    ("Recommendations", "recommendations"),
    ("Follow-up", "follow_up"),
]


# ============================================================================
# Enums and Data Classes
# ============================================================================

class DocumentType(Enum):
    """Supported document output formats."""
    PDF = "pdf"
    DOCX = "docx"
    HTML = "html"
    MD = "markdown"
    PPTX = "pptx"
    XLSX = "xlsx"


@dataclass
class DocumentTemplate:
    """Template descriptor for document generation."""
    name: str
    doc_type: DocumentType
    title: str
    sections: List[Dict[str, Any]] = field(default_factory=list)
    language: str = "he"
    author: str = "JARVIS"
    created_at: str = ""  # ISO timestamp
    metadata: Dict[str, Any] = field(default_factory=dict)

    def __post_init__(self):
        if not self.created_at:
            self.created_at = datetime.now().isoformat()

    def to_dict(self) -> Dict[str, Any]:
        return {
            "name": self.name,
            "doc_type": self.doc_type.value,
            "title": self.title,
            "sections": self.sections,
            "language": self.language,
            "author": self.author,
            "created_at": self.created_at,
            "metadata": self.metadata,
        }

    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> "DocumentTemplate":
        data = dict(data)
        if "doc_type" in data and isinstance(data["doc_type"], str):
            data["doc_type"] = DocumentType(data["doc_type"])
        return cls(**{k: v for k, v in data.items() if k in cls.__dataclass_fields__})


# ============================================================================
# DocumentGenerator
# ============================================================================

class DocumentGenerator:
    """
    Generates professional documents: PDF, Word, HTML, presentations.
    Templates for legal, medical, engineering, business documents.
    Hebrew/English support. 100% local.
    """

    def __init__(self, output_dir: str = "~/.jarvis/documents"):
        self.output_dir = Path(output_dir).expanduser()
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.templates: Dict[str, DocumentTemplate] = {}
        self._load_builtin_templates()
        self._register_hebrew_fonts()
        logger.info("DocumentGenerator initialized. Output: %s", self.output_dir)

    # -- Font helpers --------------------------------------------------------

    def _register_hebrew_fonts(self) -> None:
        """Register Hebrew-capable fonts with ReportLab if available."""
        if not HAS_REPORTLAB:
            return
        font_paths = [
            ("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", "DejaVuSans"),
            ("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", "DejaVuSans-Bold"),
            ("/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf", "LiberationSans"),
            ("/usr/share/fonts/truetype/freefont/FreeSans.ttf", "FreeSans"),
        ]
        for path, name in font_paths:
            if os.path.isfile(path):
                try:
                    pdfmetrics.registerFont(TTFont(name, path))
                    logger.debug("Registered font: %s", name)
                except Exception as exc:
                    logger.warning("Failed to register font %s: %s", name, exc)

    def _heb(self, text: str) -> str:
        """Wrap Hebrew text for ReportLab RTL handling."""
        if not text:
            return ""
        # Simple heuristic: if contains Hebrew characters, mark for RTL
        if any("\u0590" <= ch <= "\u05FF" for ch in text):
            return f'<para direction="rtl">{text}</para>'
        return text

    def _is_hebrew(self, text: str) -> bool:
        return any("\u0590" <= ch <= "\u05FF" for ch in text)

    # -- Template management -------------------------------------------------

    def _load_builtin_templates(self) -> None:
        """Populate built-in templates dictionary."""
        now = datetime.now().isoformat()

        # Hebrew legal contract
        self.templates["legal_contract_he"] = DocumentTemplate(
            name="legal_contract_he",
            doc_type=DocumentType.PDF,
            title="חוזה משפטי",
            language="he",
            sections=[
                {"heading": "הקדמה", "content": "הקדמה לחוזה בין הצדדים."},
                {"heading": "הצדדים", "content": "פרטי הצדדים החתומים על חוזה זה."},
                {"heading": "תנאים והוראות", "content": "התנאים וההוראות החלים על החוזה."},
                {"heading": "סודיות", "content": "הוראות סודיות והגבלות שימוש במידע."},
                {"heading": "סיום", "content": "תנאי סיום והפסקת החוזה."},
                {"heading": "דין חל", "content": "הדין החל על הסכם זה."},
                {"heading": "חתימות", "content": "חתימות הצדדים."},
            ],
            created_at=now,
        )

        # English legal contract
        self.templates["legal_contract_en"] = DocumentTemplate(
            name="legal_contract_en",
            doc_type=DocumentType.PDF,
            title="Legal Contract",
            language="en",
            sections=[
                {"heading": "Preamble", "content": "Preamble to the agreement between the parties."},
                {"heading": "Parties", "content": "Details of the parties signing this agreement."},
                {"heading": "Terms and Conditions", "content": "Terms and conditions governing this agreement."},
                {"heading": "Confidentiality", "content": "Confidentiality provisions and information use restrictions."},
                {"heading": "Termination", "content": "Terms of termination and contract cessation."},
                {"heading": "Governing Law", "content": "The governing law applicable to this agreement."},
                {"heading": "Signatures", "content": "Party signatures."},
            ],
            created_at=now,
        )

        # Medical report
        self.templates["medical_report"] = DocumentTemplate(
            name="medical_report",
            doc_type=DocumentType.PDF,
            title="Medical Report / דוח רפואי",
            language="en",
            sections=[
                {"heading": "Patient Information", "content": "Patient demographic details."},
                {"heading": "Referral Reason", "content": "Reason for medical examination."},
                {"heading": "Chief Complaints", "content": "Primary symptoms reported by the patient."},
                {"heading": "Physical Examination", "content": "Findings from physical examination."},
                {"heading": "Diagnostic Findings", "content": "Results of diagnostic tests and imaging."},
                {"heading": "Diagnosis", "content": "Medical diagnosis and assessment."},
                {"heading": "Recommendations", "content": "Treatment recommendations and prescriptions."},
                {"heading": "Follow-up", "content": "Follow-up plan and scheduled appointments."},
            ],
            created_at=now,
        )

        # Technical specification
        self.templates["technical_spec"] = DocumentTemplate(
            name="technical_spec",
            doc_type=DocumentType.PDF,
            title="Technical Specification",
            language="en",
            sections=[
                {"heading": "Introduction", "content": "Overview and purpose of the technical specification."},
                {"heading": "Scope", "content": "Project scope and boundaries."},
                {"heading": "System Architecture", "content": "High-level system architecture and design."},
                {"heading": "Requirements", "content": "Functional and non-functional requirements."},
                {"heading": "Interface Specifications", "content": "API and user interface specifications."},
                {"heading": "Performance Criteria", "content": "Performance benchmarks and criteria."},
                {"heading": "Security Requirements", "content": "Security and compliance requirements."},
                {"heading": "Testing Plan", "content": "Testing strategy and acceptance criteria."},
                {"heading": "Timeline", "content": "Project milestones and delivery schedule."},
            ],
            created_at=now,
        )

        # Business plan
        self.templates["business_plan"] = DocumentTemplate(
            name="business_plan",
            doc_type=DocumentType.PDF,
            title="Business Plan",
            language="en",
            sections=[
                {"heading": "Executive Summary", "content": "High-level overview of the business plan."},
                {"heading": "Company Description", "content": "Company background, mission, and vision."},
                {"heading": "Market Analysis", "content": "Industry overview, target market, and competitive landscape."},
                {"heading": "Organization & Management", "content": "Organizational structure and key personnel."},
                {"heading": "Products or Services", "content": "Description of products or services offered."},
                {"heading": "Marketing Strategy", "content": "Marketing and sales strategy."},
                {"heading": "Financial Projections", "content": "Revenue forecasts, budgets, and financial analysis."},
                {"heading": "Funding Request", "content": "Capital requirements and funding strategy."},
                {"heading": "Appendix", "content": "Supporting documents and additional information."},
            ],
            created_at=now,
        )

        # Meeting minutes
        self.templates["meeting_minutes"] = DocumentTemplate(
            name="meeting_minutes",
            doc_type=DocumentType.DOCX,
            title="Meeting Minutes",
            language="en",
            sections=[
                {"heading": "Meeting Details", "content": "Date, time, location, and participants."},
                {"heading": "Agenda", "content": "List of agenda items discussed."},
                {"heading": "Discussion Summary", "content": "Summary of key discussion points."},
                {"heading": "Decisions Made", "content": "Formal decisions and resolutions."},
                {"heading": "Action Items", "content": "Assigned tasks with owners and deadlines."},
                {"heading": "Next Steps", "content": "Planned follow-up actions and next meeting."},
            ],
            created_at=now,
        )

        # Project proposal
        self.templates["project_proposal"] = DocumentTemplate(
            name="project_proposal",
            doc_type=DocumentType.PDF,
            title="Project Proposal",
            language="en",
            sections=[
                {"heading": "Executive Summary", "content": "Brief overview of the proposed project."},
                {"heading": "Project Background", "content": "Context and rationale for the project."},
                {"heading": "Objectives", "content": "Specific, measurable project objectives."},
                {"heading": "Scope of Work", "content": "Detailed scope and deliverables."},
                {"heading": "Methodology", "content": "Approach and methodology for execution."},
                {"heading": "Timeline", "content": "Project schedule and key milestones."},
                {"heading": "Budget", "content": "Cost breakdown and budget allocation."},
                {"heading": "Team", "content": "Project team roles and responsibilities."},
                {"heading": "Risk Assessment", "content": "Identified risks and mitigation strategies."},
                {"heading": "Expected Outcomes", "content": "Anticipated results and success metrics."},
            ],
            created_at=now,
        )

    def load_template(self, template_name: str) -> DocumentTemplate:
        """Load a built-in or custom template by name."""
        if template_name in self.templates:
            # Return a shallow copy to allow per-use mutation
            tmpl = self.templates[template_name]
            return DocumentTemplate(
                name=tmpl.name,
                doc_type=tmpl.doc_type,
                title=tmpl.title,
                sections=list(tmpl.sections),
                language=tmpl.language,
                author=tmpl.author,
                created_at=tmpl.created_at,
                metadata=dict(tmpl.metadata),
            )
        tmpl_path = self.output_dir / f"{template_name}.json"
        if tmpl_path.exists():
            import json
            return DocumentTemplate.from_dict(json.loads(tmpl_path.read_text(encoding="utf-8")))
        raise KeyError(f"Template '{template_name}' not found.")

    def list_templates(self) -> List[str]:
        """Return list of all available template names."""
        custom = []
        for p in self.output_dir.glob("*.json"):
            custom.append(p.stem)
        return sorted(set(self.templates.keys()) | set(custom))

    def save_template(self, template: DocumentTemplate) -> str:
        """Save a custom template to disk for reuse."""
        import json
        path = self.output_dir / f"{template.name}.json"
        path.write_text(json.dumps(template.to_dict(), ensure_ascii=False, indent=2), encoding="utf-8")
        self.templates[template.name] = template
        return str(path)

    # -- Main entry point ----------------------------------------------------

    def create_document(self, template: DocumentTemplate) -> str:
        """Generate a document from a DocumentTemplate. Returns file path."""
        if template.doc_type == DocumentType.PDF:
            sections = template.sections or []
            content = "\n\n".join(
                f"{s.get('heading', '')}\n{s.get('content', '')}" for s in sections
            )
            return self.create_pdf(template.title, content, template.sections)
        elif template.doc_type == DocumentType.DOCX:
            content = "\n\n".join(
                f"{s.get('heading', '')}\n{s.get('content', '')}" for s in (template.sections or [])
            )
            return self.create_word(template.title, content)
        elif template.doc_type == DocumentType.HTML:
            return self._create_html(template.title, template.sections, template.language)
        elif template.doc_type == DocumentType.MD:
            return self._create_markdown(template.title, template.sections, template.language)
        elif template.doc_type == DocumentType.PPTX:
            slides = [{"title": s.get("heading", ""), "content": s.get("content", "")} for s in (template.sections or [])]
            return self.create_presentation(template.title, slides)
        elif template.doc_type == DocumentType.XLSX:
            sheets = {template.title: [[s.get("heading", ""), s.get("content", "")] for s in (template.sections or [])]}
            return self.create_spreadsheet(template.title, sheets)
        else:
            raise ValueError(f"Unsupported document type: {template.doc_type}")

    # -- PDF generation ------------------------------------------------------

    def create_pdf(self, title: str, content: str, sections: List[Dict] = None) -> str:
        """Create a PDF with proper formatting and Hebrew support."""
        safe_title = self._safe_filename(title)
        out_path = self.output_dir / f"{safe_title}_{self._timestamp()}.pdf"

        if HAS_REPORTLAB:
            self._create_pdf_reportlab(str(out_path), title, content, sections)
        elif HAS_FPDF:
            self._create_pdf_fpdf(str(out_path), title, content, sections)
        else:
            self._create_pdf_html_fallback(str(out_path), title, content, sections)

        logger.info("PDF created: %s", out_path)
        return str(out_path)

    def _create_pdf_reportlab(self, path: str, title: str, content: str, sections: List[Dict] = None) -> None:
        """Generate PDF using ReportLab with professional styling."""
        doc = SimpleDocTemplate(
            path,
            pagesize=A4,
            rightMargin=2 * cm,
            leftMargin=2 * cm,
            topMargin=2.5 * cm,
            bottomMargin=2 * cm,
        )
        story: List[Any] = []
        styles = getSampleStyleSheet()

        # Custom styles
        title_style = ParagraphStyle(
            "CustomTitle",
            parent=styles["Title"],
            fontName="DejaVuSans" if "DejaVuSans" in pdfmetrics.getRegisteredFontNames() else "Helvetica-Bold",
            fontSize=22,
            leading=28,
            alignment=TA_CENTER,
            spaceAfter=20,
        )
        heading_style = ParagraphStyle(
            "CustomHeading",
            parent=styles["Heading2"],
            fontName="DejaVuSans-Bold" if "DejaVuSans-Bold" in pdfmetrics.getRegisteredFontNames() else "Helvetica-Bold",
            fontSize=14,
            leading=18,
            spaceAfter=8,
            spaceBefore=14,
            textColor=rl_colors.HexColor("#2c3e50"),
        )
        body_style = ParagraphStyle(
            "CustomBody",
            parent=styles["Normal"],
            fontName="DejaVuSans" if "DejaVuSans" in pdfmetrics.getRegisteredFontNames() else "Helvetica",
            fontSize=10,
            leading=14,
            alignment=TA_RIGHT if self._is_hebrew(content) else TA_LEFT,
        )

        # Header line
        story.append(Spacer(1, 0.5 * cm))
        story.append(Paragraph(self._heb(title), title_style))
        story.append(Spacer(1, 0.3 * cm))

        # Date line
        date_str = datetime.now().strftime("%Y-%m-%d %H:%M")
        date_style = ParagraphStyle(
            "DateStyle",
            parent=styles["Normal"],
            fontSize=9,
            alignment=TA_CENTER,
            textColor=rl_colors.grey,
        )
        story.append(Paragraph(date_str, date_style))
        story.append(Spacer(1, 0.8 * cm))

        # Divider line
        story.append(Table([[""]], colWidths=[16 * cm], style=TableStyle([
            ("LINEBELOW", (0, 0), (-1, 0), 2, rl_colors.HexColor("#3498db")),
        ])))
        story.append(Spacer(1, 0.5 * cm))

        # Content
        if sections:
            for section in sections:
                heading = section.get("heading", "")
                body = section.get("content", "")
                if heading:
                    story.append(Paragraph(self._heb(heading), heading_style))
                if body:
                    for line in body.split("\n"):
                        if line.strip():
                            story.append(Paragraph(self._heb(line.strip()), body_style))
                story.append(Spacer(1, 0.3 * cm))
        else:
            for line in content.split("\n"):
                if line.strip():
                    story.append(Paragraph(self._heb(line.strip()), body_style))

        # Footer metadata table
        story.append(Spacer(1, 1 * cm))
        meta_data = [["Generated by JARVIS", f"Date: {date_str}", f"Page {{}} of {{}}"]]
        meta_table = Table(meta_data, colWidths=[5 * cm, 5 * cm, 6 * cm])
        meta_table.setStyle(TableStyle([
            ("FONTNAME", (0, 0), (-1, -1), "Helvetica-Oblique"),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("TEXTCOLOR", (0, 0), (-1, -1), rl_colors.grey),
            ("ALIGN", (0, 0), (-1, -1), "CENTER"),
            ("LINEBELOW", (0, 0), (-1, 0), 1, rl_colors.HexColor("#ecf0f1")),
        ]))

        doc.build(story)

    def _create_pdf_fpdf(self, path: str, title: str, content: str, sections: List[Dict] = None) -> None:
        """Generate PDF using FPDF as a fallback."""
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)

        # Title
        pdf.set_font("Arial", "B", 16)
        pdf.cell(0, 10, title.encode("latin-1", "replace").decode("latin-1"), ln=True, align="C")
        pdf.ln(5)

        # Date
        pdf.set_font("Arial", "", 9)
        date_str = datetime.now().strftime("%Y-%m-%d %H:%M")
        pdf.cell(0, 6, date_str, ln=True, align="C")
        pdf.ln(5)

        # Content
        pdf.set_font("Arial", "", 11)
        if sections:
            for section in sections:
                heading = section.get("heading", "")
                body = section.get("content", "")
                if heading:
                    pdf.set_font("Arial", "B", 13)
                    pdf.cell(0, 8, heading.encode("latin-1", "replace").decode("latin-1"), ln=True)
                if body:
                    pdf.set_font("Arial", "", 11)
                    for line in body.split("\n"):
                        safe = line.strip().encode("latin-1", "replace").decode("latin-1")
                        pdf.multi_cell(0, 6, safe)
                pdf.ln(3)
        else:
            for line in content.split("\n"):
                safe = line.strip().encode("latin-1", "replace").decode("latin-1")
                pdf.multi_cell(0, 6, safe)

        pdf.output(path)

    def _create_pdf_html_fallback(self, path: str, title: str, content: str, sections: List[Dict] = None) -> None:
        """Generate an HTML file as ultimate fallback when no PDF library is available."""
        html_path = path.replace(".pdf", ".html")
        self._render_html_file(html_path, title, content, sections)
        # Try converting via weasyprint or wkhtmltopdf if available
        try:
            subprocess.run(["weasyprint", html_path, path], check=True, capture_output=True, timeout=30)
        except Exception:
            try:
                subprocess.run(["wkhtmltopdf", "--quiet", html_path, path], check=True, capture_output=True, timeout=30)
            except Exception:
                pass  # Keep the HTML as fallback

    # -- Word (.docx) generation ---------------------------------------------

    def create_word(self, title: str, content: str) -> str:
        """Create a .docx file with python-docx."""
        safe_title = self._safe_filename(title)
        out_path = self.output_dir / f"{safe_title}_{self._timestamp()}.docx"

        if HAS_PYTHON_DOCX:
            self._create_word_docx(str(out_path), title, content)
        else:
            # Fallback to markdown
            md_path = str(out_path).replace(".docx", ".md")
            self._render_markdown_file(md_path, title, [{"heading": title, "content": content}], "en")
            return md_path

        logger.info("Word document created: %s", out_path)
        return str(out_path)

    def _create_word_docx(self, path: str, title: str, content: str) -> None:
        """Generate .docx using python-docx."""
        doc = Document()

        # Set default font
        style = doc.styles["Normal"]
        font = style.font
        font.name = "Calibri"
        font.size = Pt(11)

        # Title
        title_para = doc.add_heading(title, level=0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Date
        date_para = doc.add_paragraph()
        date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = date_para.add_run(datetime.now().strftime("%Y-%m-%d %H:%M"))
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(128, 128, 128)

        doc.add_paragraph()  # Spacer

        # Content
        for line in content.split("\n"):
            line = line.strip()
            if not line:
                continue
            if line.startswith("# ") or line.startswith("## ") or line.startswith("### "):
                level = min(line.count("#"), 3)
                text = line.lstrip("# ").strip()
                doc.add_heading(text, level=level)
            elif line.startswith("- ") or line.startswith("* "):
                doc.add_paragraph(line[2:], style="List Bullet")
            elif line[0:1].isdigit() and ". " in line[:4]:
                doc.add_paragraph(line, style="List Number")
            else:
                doc.add_paragraph(line)

        # Footer
        doc.add_paragraph()
        footer_para = doc.add_paragraph()
        footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = footer_para.add_run("Generated by JARVIS Document Generator")
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(150, 150, 150)
        run.font.italic = True

        doc.save(path)

    # -- Presentation (.pptx) generation -------------------------------------

    def create_presentation(self, title: str, slides: List[Dict]) -> str:
        """Create a .pptx presentation. Each slide: {title, content, layout}."""
        safe_title = self._safe_filename(title)
        out_path = self.output_dir / f"{safe_title}_{self._timestamp()}.pptx"

        if HAS_PYTHON_PPTX:
            self._create_pptx(str(out_path), title, slides)
        else:
            # Fallback: generate markdown outline
            md_path = str(out_path).replace(".pptx", ".md")
            self._render_presentation_as_markdown(md_path, title, slides)
            return md_path

        logger.info("Presentation created: %s", out_path)
        return str(out_path)

    def _create_pptx(self, path: str, title: str, slides: List[Dict]) -> None:
        """Generate .pptx using python-pptx."""
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)

        # Title slide
        blank_layout = prs.slide_layouts[6]  # Blank
        title_slide = prs.slides.add_slide(blank_layout)

        # Title text box
        title_box = title_slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(2.5), PptxInches(12.3), PptxInches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sub_box = title_slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(4.2), PptxInches(12.3), PptxInches(0.8)
        )
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"Generated by JARVIS | {datetime.now().strftime('%Y-%m-%d')}"
        p2.font.size = Pt(16)
        p2.font.color.rgb = PptxRGBColor(100, 100, 100)
        p2.alignment = PP_ALIGN.CENTER

        # Content slides
        for slide_data in slides:
            slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(slide_layout)

            slide_title = slide_data.get("title", "")
            slide_content = slide_data.get("content", "")
            layout_type = slide_data.get("layout", "content")

            # Slide title bar (colored background)
            title_shape = slide.shapes.add_shape(
                1,  # Rectangle
                PptxInches(0), PptxInches(0),
                PptxInches(13.333), PptxInches(1.2)
            )
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = PptxRGBColor(44, 62, 80)
            title_shape.line.fill.background()

            # Title text
            title_box = slide.shapes.add_textbox(
                PptxInches(0.5), PptxInches(0.25), PptxInches(12), PptxInches(0.8)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_title
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = PptxRGBColor(255, 255, 255)

            if layout_type == "two_column":
                # Two-column layout
                left_box = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(1.5), PptxInches(6), PptxInches(5.5)
                )
                right_box = slide.shapes.add_textbox(
                    PptxInches(6.8), PptxInches(1.5), PptxInches(6), PptxInches(5.5)
                )
                lines = slide_content.split("\n")
                mid = len(lines) // 2
                left_text = "\n".join(lines[:mid])
                right_text = "\n".join(lines[mid:])
                for box, text in [(left_box, left_text), (right_box, right_text)]:
                    tf = box.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = text
                    p.font.size = Pt(18)
                    p.line_spacing = 1.5
            else:
                # Standard content layout
                content_box = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(1.5), PptxInches(12.3), PptxInches(5.5)
                )
                tf = content_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_content
                p.font.size = Pt(20)
                p.line_spacing = 1.5

        prs.save(path)

    # -- Spreadsheet (.xlsx) generation --------------------------------------

    def create_spreadsheet(self, title: str, sheets: Dict[str, List[List]]) -> str:
        """Create an .xlsx file with multiple sheets."""
        safe_title = self._safe_filename(title)
        out_path = self.output_dir / f"{safe_title}_{self._timestamp()}.xlsx"

        if HAS_OPENPYXL:
            self._create_xlsx(str(out_path), title, sheets)
        else:
            csv_path = str(out_path).replace(".xlsx", ".csv")
            self._render_csv_fallback(csv_path, sheets)
            return csv_path

        logger.info("Spreadsheet created: %s", out_path)
        return str(out_path)

    def _create_xlsx(self, path: str, title: str, sheets: Dict[str, List[List]]) -> None:
        """Generate .xlsx using openpyxl with professional styling."""
        wb = openpyxl.Workbook()
        # Remove default sheet
        wb.remove(wb.active)

        # Styles
        header_font = Font(name="Calibri", size=12, bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="2C3E50", end_color="2C3E50", fill_type="solid")
        header_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell_border = Border(
            left=Side(style="thin", color="B0B0B0"),
            right=Side(style="thin", color="B0B0B0"),
            top=Side(style="thin", color="B0B0B0"),
            bottom=Side(style="thin", color="B0B0B0"),
        )
        alt_fill = PatternFill(start_color="ECF0F1", end_color="ECF0F1", fill_type="solid")

        for sheet_name, rows in sheets.items():
            ws = wb.create_sheet(title=sheet_name[:31])  # Excel sheet name limit

            for row_idx, row in enumerate(rows, start=1):
                for col_idx, value in enumerate(row, start=1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)
                    cell.border = cell_border
                    cell.alignment = Alignment(vertical="center", wrap_text=True)

                    if row_idx == 1:
                        cell.font = header_font
                        cell.fill = header_fill
                        cell.alignment = header_align
                    elif row_idx % 2 == 0:
                        cell.fill = alt_fill

            # Auto-size columns (within reason)
            for col_idx in range(1, 20):
                ws.column_dimensions[openpyxl.utils.get_column_letter(col_idx)].width = 18

        wb.save(path)

    # -- HTML generation -----------------------------------------------------

    def _create_html(self, title: str, sections: List[Dict], language: str = "en") -> str:
        """Create a standalone HTML document."""
        safe_title = self._safe_filename(title)
        out_path = self.output_dir / f"{safe_title}_{self._timestamp()}.html"
        self._render_html_file(str(out_path), title, None, sections, language)
        logger.info("HTML document created: %s", out_path)
        return str(out_path)

    def _render_html_file(self, path: str, title: str, content: str = None,
                          sections: List[Dict] = None, language: str = "en") -> None:
        """Render HTML file with professional styling."""
        is_rtl = language == "he"
        dir_attr = "rtl" if is_rtl else "ltr"
        align = "right" if is_rtl else "left"

        html_parts = [
            "<!DOCTYPE html>",
            f'<html lang="{language}" dir="{dir_attr}">',
            "<head>",
            '  <meta charset="UTF-8">',
            '  <meta name="viewport" content="width=device-width, initial-scale=1.0">',
            f"  <title>{title}</title>",
            "  <style>",
            "    * { box-sizing: border-box; }",
            "    body { font-family: 'Segoe UI', Arial, sans-serif; margin: 0; padding: 0; background: #f8f9fa; color: #2c3e50; }",
            "    .container { max-width: 900px; margin: 40px auto; background: #fff; padding: 50px; box-shadow: 0 2px 20px rgba(0,0,0,0.08); border-radius: 8px; }",
            "    h1 { color: #2c3e50; font-size: 28px; border-bottom: 3px solid #3498db; padding-bottom: 12px; margin-bottom: 8px; text-align: center; }",
            "    .date { text-align: center; color: #7f8c8d; font-size: 13px; margin-bottom: 30px; }",
            "    h2 { color: #2c3e50; font-size: 18px; margin-top: 28px; border-left: 4px solid #3498db; padding-left: 12px; }",
            "    p { line-height: 1.7; font-size: 14px; }",
            "    .section { margin-bottom: 24px; }",
            "    .footer { margin-top: 40px; padding-top: 20px; border-top: 1px solid #ecf0f1; text-align: center; font-size: 12px; color: #95a5a6; }",
            "    table { width: 100%; border-collapse: collapse; margin: 16px 0; }",
            "    th, td { border: 1px solid #ddd; padding: 10px; text-align: left; }",
            "    th { background: #2c3e50; color: #fff; }",
            "    tr:nth-child(even) { background: #f2f2f2; }",
            "  </style>",
            "</head>",
            "<body>",
            '  <div class="container">',
            f"    <h1>{title}</h1>",
            f"    <div class='date'>{datetime.now().strftime('%Y-%m-%d %H:%M')}</div>",
        ]

        if sections:
            for section in sections:
                heading = section.get("heading", "")
                body = section.get("content", "")
                html_parts.append(f'    <div class="section">')
                if heading:
                    html_parts.append(f"      <h2>{heading}</h2>")
                if body:
                    for line in body.split("\n"):
                        if line.strip():
                            html_parts.append(f"      <p>{line.strip()}</p>")
                html_parts.append(f"    </div>")
        elif content:
            for line in content.split("\n"):
                if line.strip():
                    html_parts.append(f"    <p>{line.strip()}</p>")

        html_parts.extend([
            '    <div class="footer">Generated by JARVIS Document Generator</div>',
            "  </div>",
            "</body>",
            "</html>",
        ])

        Path(path).write_text("\n".join(html_parts), encoding="utf-8")

    # -- Markdown generation -------------------------------------------------

    def _create_markdown(self, title: str, sections: List[Dict], language: str = "en") -> str:
        """Create a Markdown document."""
        safe_title = self._safe_filename(title)
        out_path = self.output_dir / f"{safe_title}_{self._timestamp()}.md"
        self._render_markdown_file(str(out_path), title, sections, language)
        logger.info("Markdown document created: %s", out_path)
        return str(out_path)

    def _render_markdown_file(self, path: str, title: str, sections: List[Dict], language: str = "en") -> None:
        """Render Markdown file."""
        lines = [
            f"# {title}",
            "",
            f"*Generated by JARVIS | {datetime.now().strftime('%Y-%m-%d %H:%M')}*",
            "",
            "---",
            "",
        ]
        if sections:
            for section in sections:
                heading = section.get("heading", "")
                body = section.get("content", "")
                if heading:
                    lines.append(f"## {heading}")
                    lines.append("")
                if body:
                    lines.extend(body.split("\n"))
                    lines.append("")
        Path(path).write_text("\n".join(lines), encoding="utf-8")

    def _render_presentation_as_markdown(self, path: str, title: str, slides: List[Dict]) -> None:
        """Render presentation outline as markdown fallback."""
        lines = [f"# {title}", "", f"*Generated by JARVIS | {datetime.now().strftime('%Y-%m-%d')}*", "", "---", ""]
        for i, slide in enumerate(slides, 1):
            lines.append(f"## Slide {i}: {slide.get('title', '')}")
            lines.append("")
            lines.append(slide.get("content", ""))
            lines.append("")
            lines.append("---")
            lines.append("")
        Path(path).write_text("\n".join(lines), encoding="utf-8")

    def _render_csv_fallback(self, path: str, sheets: Dict[str, List[List]]) -> None:
        """Render first sheet as CSV fallback."""
        import csv
        first_sheet = next(iter(sheets.values()), [])
        with open(path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerows(first_sheet)

    # -- Specialized document generators -------------------------------------

    def create_legal_contract(self, contract_type: str, parties: List[str], terms: List[str]) -> str:
        """Generate a legal contract. Types: nda, employment, service, partnership, license."""
        contract_type = contract_type.lower()
        type_titles = {
            "nda": "Non-Disclosure Agreement",
            "employment": "Employment Agreement",
            "service": "Service Agreement",
            "partnership": "Partnership Agreement",
            "license": "Software License Agreement",
        }
        title = type_titles.get(contract_type, f"{contract_type.title()} Agreement")
        safe_title = self._safe_filename(title)
        out_path = self.output_dir / f"{safe_title}_{self._timestamp()}.pdf"

        # Build contract content
        sections = self._build_legal_sections(contract_type, parties, terms)

        # Generate PDF
        content = "\n\n".join(f"{s['heading']}\n{s['content']}" for s in sections)
        if HAS_REPORTLAB:
            self._create_pdf_reportlab(str(out_path), title, content, sections)
        elif HAS_FPDF:
            self._create_pdf_fpdf(str(out_path), title, content, sections)
        else:
            self._create_pdf_html_fallback(str(out_path), title, content, sections)

        logger.info("Legal contract created: %s", out_path)
        return str(out_path)

    def _build_legal_sections(self, contract_type: str, parties: List[str], terms: List[str]) -> List[Dict]:
        """Build legal contract section content."""
        sections = []
        date_str = datetime.now().strftime("%B %d, %Y")

        # Preamble
        party_str = ", ".join(parties) if parties else "[PARTY NAMES]"
        preamble = (
            f"This {contract_type.replace('_', ' ').title()} Agreement (\"Agreement\") is entered into "
            f"as of {date_str}, by and between {party_str}."
        )
        sections.append({"heading": "Preamble", "content": preamble})

        # Parties
        party_details = "\n".join(f"Party {i+1}: {name}" for i, name in enumerate(parties))
        sections.append({"heading": "Parties", "content": party_details or "[PARTY DETAILS TO BE FILLED IN]"})

        # Recitals
        sections.append({"heading": "Recitals", "content": (
            "WHEREAS, the Parties desire to set forth the terms and conditions "
            "governing their relationship as specified herein;\n\n"
            "NOW, THEREFORE, in consideration of the mutual covenants and agreements "
            "contained herein, the Parties agree as follows:"
        )})

        # Terms
        sections.append({"heading": "Terms and Conditions", "content": "\n".join(
            f"{i+1}. {term}" for i, term in enumerate(terms)
        ) or "[TERMS TO BE SPECIFIED]"})

        # Specific clauses per contract type
        if contract_type == "nda":
            sections.append({"heading": "Confidentiality Obligations", "content": (
                "Each Party agrees to:\n"
                "(a) Hold all Confidential Information in strict confidence;\n"
                "(b) Not disclose Confidential Information to any third parties;\n"
                "(c) Use Confidential Information solely for the purpose of this Agreement;\n"
                "(d) Return or destroy all Confidential Information upon termination."
            )})
        elif contract_type == "employment":
            sections.append({"heading": "Position and Duties", "content": (
                "The Employee agrees to perform the duties of the position as described "
                "by the Employer, including but not limited to responsibilities assigned "
                "from time to time."
            )})
            sections.append({"heading": "Compensation", "content": (
                "The Employer shall pay the Employee a salary as agreed upon, "
                "subject to applicable tax withholdings and deductions."
            )})
        elif contract_type == "service":
            sections.append({"heading": "Scope of Services", "content": (
                "The Service Provider agrees to provide the services described herein "
                "in a professional and workmanlike manner, in accordance with industry standards."
            )})
        elif contract_type == "partnership":
            sections.append({"heading": "Capital Contributions", "content": (
                "Each Partner shall contribute capital as specified in Schedule A. "
                "Profits and losses shall be distributed in proportion to capital contributions "
                "unless otherwise agreed."
            )})
        elif contract_type == "license":
            sections.append({"heading": "License Grant", "content": (
                "Licensor hereby grants Licensee a non-exclusive, non-transferable license "
                "to use the Software subject to the terms and conditions of this Agreement."
            )})

        # Standard closing sections
        sections.append({"heading": "Term and Termination", "content": (
            "This Agreement shall commence on the Effective Date and continue until terminated "
            "by either Party with [NOTICE PERIOD] written notice, or as otherwise specified herein."
        )})

        sections.append({"heading": "Governing Law", "content": (
            "This Agreement shall be governed by and construed in accordance with the laws "
            "of the jurisdiction in which the primary place of business of the first Party is located, "
            "without regard to conflict of law principles."
        )})

        sections.append({"heading": "Limitation of Liability", "content": (
            "In no event shall either Party be liable for any indirect, incidental, special, "
            "consequential, or punitive damages arising out of or related to this Agreement."
        )})

        sections.append({"heading": "Signatures", "content": (
            "IN WITNESS WHEREOF, the Parties have executed this Agreement as of the date first written above.\n\n"
            + "\n\n".join(f"_________________________\nParty {i+1}: {name}\nDate: _______________"
                         for i, name in enumerate(parties))
            + "\n\n_________________________\nWitness: _______________\nDate: _______________"
        )})

        return sections

    def create_medical_report(self, patient_info: Dict, findings: List[str]) -> str:
        """Generate a medical report with professional format."""
        title = f"Medical Report - {patient_info.get('name', 'Unknown Patient')}"
        safe_title = self._safe_filename(title)
        out_path = self.output_dir / f"{safe_title}_{self._timestamp()}.pdf"

        lang = patient_info.get("language", "en")
        sections = self._build_medical_sections(patient_info, findings, lang)

        content = "\n\n".join(f"{s['heading']}\n{s['content']}" for s in sections)

        if HAS_REPORTLAB:
            self._create_pdf_reportlab(str(out_path), title, content, sections)
        elif HAS_FPDF:
            self._create_pdf_fpdf(str(out_path), title, content, sections)
        else:
            self._create_pdf_html_fallback(str(out_path), title, content, sections)

        logger.info("Medical report created: %s", out_path)
        return str(out_path)

    def _build_medical_sections(self, patient_info: Dict, findings: List[str], lang: str) -> List[Dict]:
        """Build medical report sections."""
        sections = []
        med_sections = MEDICAL_SECTIONS_HE if lang == "he" else MEDICAL_SECTIONS_EN

        # Patient info section
        info_lines = []
        for key, value in patient_info.items():
            if key != "language":
                info_lines.append(f"{key.replace('_', ' ').title()}: {value}")
        sections.append({"heading": med_sections[0][0], "content": "\n".join(info_lines)})

        # Findings
        findings_text = "\n".join(f"- {finding}" for finding in findings) if findings else "No findings recorded."
        sections.append({"heading": med_sections[4][0], "content": findings_text})

        # Placeholder sections for other categories
        for label, key in med_sections[1:4] + med_sections[5:]:
            sections.append({"heading": label, "content": f"[{label} - to be completed by physician]"})

        return sections

    def create_technical_spec(self, project_name: str, requirements: List[str]) -> str:
        """Generate a technical specification document."""
        title = f"Technical Specification - {project_name}"
        safe_title = self._safe_filename(title)
        out_path = self.output_dir / f"{safe_title}_{self._timestamp()}.pdf"

        sections = [
            {"heading": "1. Introduction",
             "content": f"This document specifies the technical requirements for {project_name}.\n\n"
                        "It covers system architecture, functional requirements, non-functional requirements, "
                        "and implementation guidelines."},
            {"heading": "2. Scope",
             "content": "The scope of this technical specification includes all components, interfaces, "
                        "and subsystems necessary for the successful implementation of the project."},
            {"heading": "3. System Architecture",
             "content": "[System architecture diagram and description to be inserted here.]\n\n"
                        "The system follows a modular architecture with clearly defined interfaces between components."},
            {"heading": "4. Functional Requirements",
             "content": "\n".join(f"FR-{i+1:03d}: {req}" for i, req in enumerate(requirements))
             or "[REQUIREMENTS TO BE SPECIFIED]"},
            {"heading": "5. Non-Functional Requirements",
             "content": "NFR-001: The system shall support at least 1000 concurrent users.\n"
                        "NFR-002: Response time for critical operations shall not exceed 200ms.\n"
                        "NFR-003: System availability shall be 99.9% or higher.\n"
                        "NFR-004: All data shall be encrypted at rest and in transit.\n"
                        "NFR-005: The system shall be deployable within 30 minutes."},
            {"heading": "6. Interface Specifications",
             "content": "[API endpoints, data formats, and communication protocols to be documented here.]"},
            {"heading": "7. Performance Criteria",
             "content": "[Benchmarks, load testing criteria, and performance targets to be specified.]"},
            {"heading": "8. Security Requirements",
             "content": "All components shall implement authentication and authorization mechanisms.\n"
                        "Input validation shall be performed on all external interfaces.\n"
                        "Security vulnerabilities shall be addressed within 48 hours of discovery."},
            {"heading": "9. Testing Strategy",
             "content": "Unit tests: Minimum 80% code coverage.\n"
                        "Integration tests: All API endpoints and component interactions.\n"
                        "Performance tests: Load testing under expected peak conditions.\n"
                        "Security tests: Penetration testing and vulnerability scanning."},
            {"heading": "10. Timeline and Milestones",
             "content": "[Project timeline with key milestones and deliverable dates to be inserted.]"},
        ]

        content = "\n\n".join(f"{s['heading']}\n{s['content']}" for s in sections)

        if HAS_REPORTLAB:
            self._create_pdf_reportlab(str(out_path), title, content, sections)
        elif HAS_FPDF:
            self._create_pdf_fpdf(str(out_path), title, content, sections)
        else:
            self._create_pdf_html_fallback(str(out_path), title, content, sections)

        logger.info("Technical specification created: %s", out_path)
        return str(out_path)

    def create_business_plan(self, business_name: str, sections: Dict[str, str]) -> str:
        """Generate a business plan document."""
        title = f"Business Plan - {business_name}"
        safe_title = self._safe_filename(title)
        out_path = self.output_dir / f"{safe_title}_{self._timestamp()}.pdf"

        default_sections = [
            {"heading": "Executive Summary",
             "content": sections.get("executive_summary",
                 "This business plan outlines the vision, strategy, and financial projections for "
                 f"{business_name}. The company aims to address identified market needs through innovative "
                 "products and services.")},
            {"heading": "Company Description",
             "content": sections.get("company_description",
                 f"{business_name} is a forward-thinking company dedicated to delivering exceptional "
                 "value to its customers. Our mission is to [MISSION STATEMENT].")},
            {"heading": "Market Analysis",
             "content": sections.get("market_analysis",
                 "The target market consists of [DESCRIBE TARGET MARKET]. Market size is estimated "
                 "at $[X] billion with a projected CAGR of [Y]% over the next 5 years.")},
            {"heading": "Products and Services",
             "content": sections.get("products_services",
                 "[DESCRIBE PRODUCTS/SERVICES OFFERED, UNIQUE VALUE PROPOSITION, AND COMPETITIVE ADVANTAGES.]")},
            {"heading": "Marketing Strategy",
             "content": sections.get("marketing_strategy",
                 "Our go-to-market strategy includes digital marketing, content marketing, strategic partnerships, "
                 "and direct sales. Customer acquisition cost is projected at $[X] with LTV of $[Y].")},
            {"heading": "Management Team",
             "content": sections.get("management_team",
                 "[DESCRIBE KEY MANAGEMENT PERSONNEL, THEIR BACKGROUNDS, AND ROLES.]")},
            {"heading": "Financial Projections",
             "content": sections.get("financial_projections",
                 "Year 1 Revenue: $[X] | Expenses: $[Y] | Net Income: $[Z]\n"
                 "Year 2 Revenue: $[X] | Expenses: $[Y] | Net Income: $[Z]\n"
                 "Year 3 Revenue: $[X] | Expenses: $[Y] | Net Income: $[Z]\n"
                 "Break-even projected in Month [N].")},
            {"heading": "Funding Requirements",
             "content": sections.get("funding_requirements",
                 "Total funding required: $[AMOUNT]\n"
                 "Use of funds:\n"
                 "- Product Development: [X]%\n"
                 "- Marketing and Sales: [Y]%\n"
                 "- Operations: [Z]%\n"
                 "- Working Capital: [W]%")},
            {"heading": "Risk Analysis",
             "content": sections.get("risk_analysis",
                 "Key risks include market competition, regulatory changes, technology disruption, and "
                 "economic downturns. Mitigation strategies have been developed for each identified risk.")},
            {"heading": "Appendix",
             "content": "[SUPPORTING DOCUMENTS, MARKET RESEARCH DATA, FINANCIAL STATEMENTS, AND ADDITIONAL REFERENCES.]"},
        ]

        content = "\n\n".join(f"{s['heading']}\n{s['content']}" for s in default_sections)

        if HAS_REPORTLAB:
            self._create_pdf_reportlab(str(out_path), title, content, default_sections)
        elif HAS_FPDF:
            self._create_pdf_fpdf(str(out_path), title, content, default_sections)
        else:
            self._create_pdf_html_fallback(str(out_path), title, content, default_sections)

        logger.info("Business plan created: %s", out_path)
        return str(out_path)

    # -- Utility methods -----------------------------------------------------

    @staticmethod
    def _safe_filename(name: str) -> str:
        """Sanitize a string for use as a filename."""
        keep = "abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789_-"
        return "".join(c if c in keep else "_" for c in name)[:80]

    @staticmethod
    def _timestamp() -> str:
        """Generate a compact timestamp string for filenames."""
        return datetime.now().strftime("%Y%m%d_%H%M%S")

    def get_stats(self) -> Dict[str, Any]:
        """Return generator statistics and capabilities."""
        return {
            "output_dir": str(self.output_dir),
            "templates_available": len(self.templates),
            "template_names": self.list_templates(),
            "capabilities": {
                "pdf_reportlab": HAS_REPORTLAB,
                "pdf_fpdf": HAS_FPDF,
                "docx": HAS_PYTHON_DOCX,
                "pptx": HAS_PYTHON_PPTX,
                "xlsx": HAS_OPENPYXL,
            },
        }


# ============================================================================
# MockDocumentGenerator — same interface, placeholder files
# ============================================================================

class MockDocumentGenerator(DocumentGenerator):
    """
    Mock implementation of DocumentGenerator that creates placeholder files
    with the same interface. Useful for testing and development environments
    where full generation is not required.
    """

    def __init__(self, output_dir: str = "~/.jarvis/documents"):
        self.output_dir = Path(output_dir).expanduser()
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.templates: Dict[str, DocumentTemplate] = {}
        self._load_builtin_templates()
        logger.info("MockDocumentGenerator initialized (placeholders only).")

    def create_pdf(self, title: str, content: str, sections: List[Dict] = None) -> str:
        """Create a placeholder PDF file."""
        safe_title = self._safe_filename(title)
        out_path = self.output_dir / f"{safe_title}_{self._timestamp()}.pdf"
        out_path.write_text(
            f"%PDF-1.4 PLACEHOLDER\n%\n"
            f"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n"
            f"2 0 obj\n<< /Type /Pages /Kids [] /Count 0 >>\nendobj\n"
            f"xref\ntrailer\n<< /Root 1 0 R >>\n%%EOF\n",
            encoding="latin-1",
        )
        logger.info("[MOCK] Placeholder PDF created: %s", out_path)
        return str(out_path)

    def create_word(self, title: str, content: str) -> str:
        """Create a placeholder .docx file."""
        safe_title = self._safe_filename(title)
        out_path = self.output_dir / f"{safe_title}_{self._timestamp()}.docx"
        # Minimal valid zip-based docx structure
        self._write_minimal_docx(str(out_path), title, content)
        logger.info("[MOCK] Placeholder DOCX created: %s", out_path)
        return str(out_path)

    def create_presentation(self, title: str, slides: List[Dict]) -> str:
        """Create a placeholder .pptx file."""
        safe_title = self._safe_filename(title)
        out_path = self.output_dir / f"{safe_title}_{self._timestamp()}.pptx"
        out_path.write_bytes(b"PK\x03\x04" + b"\x00" * 26)  # Minimal ZIP header
        logger.info("[MOCK] Placeholder PPTX created: %s", out_path)
        return str(out_path)

    def create_spreadsheet(self, title: str, sheets: Dict[str, List[List]]) -> str:
        """Create a placeholder .xlsx file."""
        safe_title = self._safe_filename(title)
        out_path = self.output_dir / f"{safe_title}_{self._timestamp()}.xlsx"
        out_path.write_bytes(b"PK\x03\x04" + b"\x00" * 26)  # Minimal ZIP header
        logger.info("[MOCK] Placeholder XLSX created: %s", out_path)
        return str(out_path)

    def create_legal_contract(self, contract_type: str, parties: List[str], terms: List[str]) -> str:
        """Create a placeholder legal contract PDF."""
        return self.create_pdf(
            f"{contract_type.title()} Agreement",
            f"Legal contract between {', '.join(parties)}\nTerms: {', '.join(terms[:3])}...",
        )

    def create_medical_report(self, patient_info: Dict, findings: List[str]) -> str:
        """Create a placeholder medical report PDF."""
        patient_name = patient_info.get("name", "Unknown")
        return self.create_pdf(
            f"Medical Report - {patient_name}",
            f"Patient: {patient_name}\nFindings: {', '.join(findings[:3])}...",
        )

    def create_technical_spec(self, project_name: str, requirements: List[str]) -> str:
        """Create a placeholder technical spec PDF."""
        return self.create_pdf(
            f"Technical Spec - {project_name}",
            f"Requirements:\n" + "\n".join(f"- {r}" for r in requirements[:5]),
        )

    def create_business_plan(self, business_name: str, sections: Dict[str, str]) -> str:
        """Create a placeholder business plan PDF."""
        return self.create_pdf(
            f"Business Plan - {business_name}",
            f"Business Plan for {business_name}\nSections: {', '.join(sections.keys())}",
        )

    def _write_minimal_docx(self, path: str, title: str, content: str) -> None:
        """Write a minimal valid .docx (which is a ZIP file)."""
        import zipfile
        from io import BytesIO
        buf = BytesIO()
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
            zf.writestr("[Content_Types].xml",
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
                '<Default Extension="xml" ContentType="application/xml"/>'
                '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
                '<Override PartName="/word/document.xml" '
                'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
                '</Types>')
            zf.writestr("_rels/.rels",
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                '<Relationship Id="rId1" '
                'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
                'Target="word/document.xml"/>'
                '</Relationships>')
            zf.writestr("word/_rels/document.xml.rels",
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                '</Relationships>')
            doc_xml = (
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
                '<w:body>'
                '<w:p><w:pPr><w:pStyle w:val="Title"/></w:pPr>'
                f'<w:r><w:t>{title}</w:t></w:r></w:p>'
                '<w:p><w:r><w:t>' + content.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;") + '</w:t></w:r></w:p>'
                '</w:body></w:document>'
            )
            zf.writestr("word/document.xml", doc_xml)
        Path(path).write_bytes(buf.getvalue())


# ============================================================================
# Factory function
# ============================================================================

def get_document_generator(use_mock: bool = False, output_dir: str = "~/.jarvis/documents") -> DocumentGenerator:
    """
    Factory function to get a DocumentGenerator instance.

    Args:
        use_mock: If True, returns MockDocumentGenerator (placeholder files only).
        output_dir: Directory where generated documents are saved.

    Returns:
        DocumentGenerator or MockDocumentGenerator instance.
    """
    if use_mock:
        return MockDocumentGenerator(output_dir=output_dir)
    return DocumentGenerator(output_dir=output_dir)


# ============================================================================
# Module-level convenience functions
# ============================================================================

def quick_pdf(title: str, content: str, output_dir: str = "~/.jarvis/documents") -> str:
    """Quickly generate a PDF document."""
    return get_document_generator(output_dir=output_dir).create_pdf(title, content)


def quick_docx(title: str, content: str, output_dir: str = "~/.jarvis/documents") -> str:
    """Quickly generate a Word document."""
    return get_document_generator(output_dir=output_dir).create_word(title, content)


def quick_presentation(title: str, slides: List[Dict], output_dir: str = "~/.jarvis/documents") -> str:
    """Quickly generate a presentation."""
    return get_document_generator(output_dir=output_dir).create_presentation(title, slides)


def quick_spreadsheet(title: str, sheets: Dict[str, List[List]], output_dir: str = "~/.jarvis/documents") -> str:
    """Quickly generate a spreadsheet."""
    return get_document_generator(output_dir=output_dir).create_spreadsheet(title, sheets)


# ============================================================================
# CLI / self-test
# ============================================================================

if __name__ == "__main__":
    # Self-test entry point – lightweight logger already configured
    print("=" * 60)
    print("JARVIS Document Generator - Self Test")
    print("=" * 60)

    dg = get_document_generator()
    stats = dg.get_stats()

    print(f"\nOutput directory: {stats['output_dir']}")
    print(f"Templates available: {stats['templates_available']}")
    print(f"Template names: {', '.join(stats['template_names'])}")
    print(f"\nCapabilities:")
    for cap, available in stats['capabilities'].items():
        status = "✓" if available else "✗"
        print(f"  {status} {cap}")

    # Test each document type
    print("\n--- Generating test documents ---")

    # PDF
    pdf_path = dg.create_pdf(
        "Test Document",
        "This is a test PDF document generated by JARVIS.",
        [{"heading": "Section 1", "content": "Content of section 1."},
         {"heading": "Section 2", "content": "Content of section 2."}]
    )
    print(f"  PDF: {pdf_path}")

    # Word
    word_path = dg.create_word(
        "Test Word Document",
        "This is a test Word document.\n\n## Features\n- Bullet 1\n- Bullet 2\n\nSome final text."
    )
    print(f"  DOCX: {word_path}")

    # Presentation
    pptx_path = dg.create_presentation("Test Presentation", [
        {"title": "Slide 1 - Introduction", "content": "Welcome to the presentation.", "layout": "title"},
        {"title": "Slide 2 - Overview", "content": "Key points:\n- Point A\n- Point B\n- Point C", "layout": "content"},
        {"title": "Slide 3 - Details", "content": "Left side content.\n\nRight side content.", "layout": "two_column"},
    ])
    print(f"  PPTX: {pptx_path}")

    # Spreadsheet
    xlsx_path = dg.create_spreadsheet("Test Spreadsheet", {
        "Summary": [
            ["Item", "Description", "Value"],
            ["A1", "First item", "100"],
            ["A2", "Second item", "200"],
            ["A3", "Third item", "300"],
        ],
        "Details": [
            ["Category", "Detail", "Status"],
            ["Dev", "Implementation", "Done"],
            ["Test", "Unit tests", "In Progress"],
            ["Deploy", "Production", "Pending"],
        ],
    })
    print(f"  XLSX: {xlsx_path}")

    # Legal contract
    contract_path = dg.create_legal_contract(
        "nda",
        ["Acme Corporation", "Beta Technologies Ltd."],
        ["All confidential information shall be protected for a period of 5 years.",
         "No reverse engineering of provided materials is permitted.",
         "Information may only be disclosed to employees with a need to know."]
    )
    print(f"  Legal Contract (NDA): {contract_path}")

    # Medical report
    med_path = dg.create_medical_report(
        {"name": "John Doe", "age": "45", "gender": "Male", "id": "123456",
         "date_of_birth": "1980-03-15", "language": "en"},
        ["Blood pressure: 120/80 mmHg (normal)",
         "Heart rate: 72 bpm (regular rhythm)",
         "Chest X-ray: No acute findings",
         "Blood tests: All values within normal range"]
    )
    print(f"  Medical Report: {med_path}")

    # Technical specification
    spec_path = dg.create_technical_spec(
        "Cloud Migration Project",
        ["The system shall support zero-downtime migration.",
         "All data must be encrypted using AES-256.",
         "API response time shall not exceed 100ms at p99.",
         "The system shall auto-scale based on CPU and memory metrics.",
         "Audit logs shall be retained for a minimum of 7 years."]
    )
    print(f"  Technical Spec: {spec_path}")

    # Business plan
    bp_path = dg.create_business_plan(
        "TechVenture AI",
        {"executive_summary": "AI-powered analytics platform for enterprise customers.",
         "market_analysis": "The AI analytics market is projected to reach $50B by 2028.",
         "financial_projections": "Year 1: $2M revenue | Year 2: $8M | Year 3: $25M"}
    )
    print(f"  Business Plan: {bp_path}")

    # Template-based document
    tmpl = dg.load_template("project_proposal")
    tmpl.title = "AI Infrastructure Proposal"
    tmpl_doc_path = dg.create_document(tmpl)
    print(f"  Template Document: {tmpl_doc_path}")

    print("\n" + "=" * 60)
    print("All tests completed successfully!")
    print(f"Documents saved to: {dg.output_dir}")
    print("=" * 60)
